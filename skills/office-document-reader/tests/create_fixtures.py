#!/usr/bin/env python3
"""Generate temporary DOCX, PPTX, and XLSX fixtures for office-document-reader tests."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path


def create_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Fixture DOCX Heading Alpha", level=1)
    doc.add_paragraph("Fixture DOCX paragraph beta with distinctive text.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "FixtureTableHeaderA"
    table.cell(0, 1).text = "FixtureTableHeaderB"
    table.cell(1, 0).text = "FixtureCellValue42"
    table.cell(1, 1).text = "FixtureCellValue99"
    doc.save(path)


def create_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Fixture PPTX Title Gamma"
    slide.placeholders[1].text = "Fixture PPTX body delta distinctive text."

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Fixture speaker notes epsilon distinctive."

    prs.save(path)


def create_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "FixtureSheetOne"
    ws1["A1"] = "FixtureXlsxValue111"
    ws1["B1"] = "FixtureXlsxValue222"

    ws2 = wb.create_sheet("FixtureSheetTwo")
    ws2["A1"] = "FixtureXlsxValue333"
    ws2["B1"] = "FixtureXlsxValue444"

    wb.save(path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Create office document test fixtures.")
    parser.add_argument("output_dir", type=Path, help="Directory to write fixture files into")
    args = parser.parse_args()

    out = args.output_dir
    out.mkdir(parents=True, exist_ok=True)

    create_docx(out / "fixture.docx")
    create_pptx(out / "fixture.pptx")
    create_xlsx(out / "fixture.xlsx")

    print(out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
